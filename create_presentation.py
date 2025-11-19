#!/usr/bin/env python3
"""
Create PowerPoint presentation for "Death to PowerPoint" lunch & learn
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(20, 20, 30))

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    top = Inches(4)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_color=None):
    """Add a content slide with title and bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(20, 20, 30))

    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = highlight_color or RGBColor(100, 200, 255)

    # Bullets
    top = Inches(1.8)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(20, 20, 30))

    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 200, 255)

    # Left column title
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(4.2)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 150, 100)

    # Left column bullets
    top = Inches(2.2)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(8)

    # Right column title
    left = Inches(5.2)
    top = Inches(1.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 255, 150)

    # Right column bullets
    top = Inches(2.2)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(8)

    return slide

def add_code_slide(prs, title, code_text, caption=""):
    """Add a slide with code snippet"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(20, 20, 30))

    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 200, 255)

    # Code box background
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
    shape.line.color.rgb = RGBColor(60, 64, 72)

    # Code text
    left = Inches(0.7)
    top = Inches(1.7)
    width = Inches(8.6)
    height = Inches(4.1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Caption
    if caption:
        left = Inches(0.5)
        top = Inches(6.2)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_stats_slide(prs, title, stats):
    """Add a slide with big stats"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(20, 20, 30))

    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 200, 255)

    # Stats in columns
    num_stats = len(stats)
    col_width = 9 / num_stats

    for i, (number, label) in enumerate(stats):
        left = Inches(0.5 + i * col_width)

        # Number
        top = Inches(2.5)
        width = Inches(col_width - 0.2)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 200, 100)
        p.alignment = PP_ALIGN.CENTER

        # Label
        top = Inches(4)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Death to PowerPoint",
        "How I Ship Customer Demos in Hours, Not Weeks\n\nPravin Vasudevan | ANZ Field Engineering"
    )

    # Slide 2: The Problem
    add_content_slide(
        prs,
        "The Old Way Was Killing Me",
        [
            "Customer asks for a demo on Tuesday",
            "Spend 2-3 days building POC from scratch",
            "Weekend work to polish before meeting",
            "PowerPoint decks that nobody remembers",
            "Context switching between N customers",
            "Repeated patterns coded from memory each time"
        ],
        RGBColor(255, 100, 100)
    )

    # Slide 3: The Insight
    add_content_slide(
        prs,
        "The Paradigm Shift",
        [
            "I stopped WRITING code",
            "I started ORCHESTRATING code",
            "AI agents do the typing",
            "I do the thinking and directing",
            "Show > Tell (live demos beat slides)",
            "Ship fast, iterate faster"
        ],
        RGBColor(100, 255, 150)
    )

    # Slide 4: My Stack
    add_two_column_slide(
        prs,
        "The Agentic Coding Stack",
        "Claude Code (CLI)",
        [
            "Terminal-based AI agent",
            "Reads/writes/edits files",
            "Runs tests, linting, git",
            "Great for automation tasks",
            "Batch operations",
            "CI/CD setup"
        ],
        "Cursor (IDE)",
        [
            "VS Code fork with AI",
            "Inline completions",
            "Chat with codebase context",
            "Visual debugging",
            "Feature development",
            "Refactoring"
        ]
    )

    # Slide 5: The Workflow
    add_content_slide(
        prs,
        "How They Work Together",
        [
            "Both connect to same GitHub repo",
            "Cursor: Write features (visual context helps)",
            "Claude Code: Tests, linting, docs, git ops",
            "Push/pull to sync between them",
            "Parallel work on different tasks",
            "Like having two 10x engineers"
        ]
    )

    # Slide 6: Live Demo Preview
    add_content_slide(
        prs,
        "What You're About to See",
        [
            "Battery Trading AI Assistant (real customer demo)",
            "LangGraph agent with RAG + SQL tools",
            "Cursor: Add a new agent tool",
            "Claude Code: Add pytest infrastructure",
            "Claude Code: Add pre-commit hooks",
            "All in ~15 minutes"
        ],
        RGBColor(255, 200, 100)
    )

    # Slide 7: Code Pattern - The Good Stuff
    add_code_slide(
        prs,
        "Pattern: Self-Documenting Tools",
        '''@tool
def get_battery_status(
    battery_id: Annotated[str,
        "Battery ID (RESS2, DPNTBESS) or 'all'"] = "all"
) -> str:
    """Get current SoC and capabilities for batteries."""

    # AI reads these annotations for tool selection
    # Also serves as documentation
    # Type hints enable testing
    ...''',
        "Annotated types = docs + AI context + type safety"
    )

    # Slide 8: Key Patterns
    add_content_slide(
        prs,
        "Patterns That Work",
        [
            "Read first, edit second (show existing patterns)",
            "Specific prompts > vague prompts",
            "Break big tasks into small chunks",
            "Use factory patterns for testability",
            "Centralize config, no magic strings",
            "Type hints everywhere (AI reads them)"
        ]
    )

    # Slide 9: Pitfalls
    add_content_slide(
        prs,
        "When AI Agents Go Rogue",
        [
            "Context window limits (break up big files)",
            "Hallucinated APIs (always verify against docs)",
            "Wrong patterns (be explicit: 'use pytest not unittest')",
            "Over-eager refactoring (scope your asks tightly)",
            "Inconsistent style (read existing code first)",
            "The 'delete everything and start over' impulse"
        ],
        RGBColor(255, 100, 100)
    )

    # Slide 10: Results
    add_stats_slide(
        prs,
        "The Results",
        [
            ("8+", "Customer demos\nthis quarter"),
            ("3-5 hrs", "Average build\ntime"),
            ("0", "PowerPoint\nslides"),
            ("10x", "Perceived\nproductivity")
        ]
    )

    # Slide 11: Getting Started
    add_content_slide(
        prs,
        "Try It Today",
        [
            "Claude Code: claude.ai/code (free tier available)",
            "Cursor: cursor.sh (free tier available)",
            "Start small: 'write tests for this file'",
            "Read the code first, then ask for edits",
            "Use specific prompts with context",
            "I'll share my templates and prompts after"
        ]
    )

    # Slide 12: Q&A
    add_title_slide(
        prs,
        "Questions?",
        "Let's debug together\n\n@pravinva | pravin.vasudevan@databricks.com"
    )

    # Save
    output_path = "/home/user/battery-trading-assistant/death_to_powerpoint.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
